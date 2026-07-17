"""Manager-update report: gather June (current-month) work and render a slide.

Structure: open Epics assigned to me -> their Stories -> Sub-tasks whose Target
Completion Date falls in the chosen month. Each sub-task gets a RAG status; each
Story shows a progress summary. Rendered as a single-table PowerPoint slide.
"""
from __future__ import annotations

import calendar
import io
import re
from datetime import date

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from . import fields as fields_mod
from .jira_client import JiraClient, adf_to_text
from .local_fields import LocalFieldsStore

START_DATE_FIELD = "customfield_10015"

# RAG palette (R, G, B)
GREEN = RGBColor(0x36, 0xB3, 0x7E)
AMBER = RGBColor(0xFF, 0x99, 0x1F)
RED = RGBColor(0xDE, 0x35, 0x0B)
BLUE = RGBColor(0x05, 0x29, 0xCC)
GREY = RGBColor(0x97, 0xA0, 0xAF)
HEADER_BG = RGBColor(0x05, 0x29, 0xCC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x17, 0x2B, 0x4D)


# Phrases in comments that indicate a delay / schedule risk. Multi-word where
# possible to avoid false positives (e.g. "late" inside "related").
DELAY_PHRASES = [
    "delay", "postpone", "push back", "pushed back", "push out", "pushed out",
    "slipp", "behind schedule", "falling behind", "fall behind",
    "blocked", "blocker", "blocking", "on hold", "on-hold",
    "waiting on", "waiting for", "awaiting", "reschedul",
    "moving the tcd", "move the tcd", "moved the tcd", "moving tcd",
    "revised tcd", "new tcd", "extend the timeline", "extended timeline",
    "overrun", "at risk", "running late", "held up", "hold up", "stuck",
    "revised timeline", "revised date", "new timeline", "miss the",
    "won't make", "wont make", "not on track", "off track",
]


def mentions_delay(text: str) -> bool:
    t = (text or "").lower()
    return any(p in t for p in DELAY_PHRASES)


def rag_status(name: str, category: str, end, comment_text: str = "",
               today: date | None = None):
    """Return (label, colour) RAG assessment for a sub-task.

    "At Risk" (amber) is raised only when a delay/schedule concern is mentioned
    in the comments. Overdue (past target date) stays red regardless.
    """
    today = today or date.today()
    if (name or "").lower() == "cancelled":
        return "Cancelled", GREY
    if category == "done":
        return "Done", BLUE
    wd = fields_mod.working_days_until(end, today)
    if wd is not None and wd < 0:
        return "Overdue", RED
    if mentions_delay(comment_text):
        return "At Risk", AMBER
    if wd is None:
        return "No date", GREY
    return "On Track", GREEN


def _one_line(text: str, limit: int = 160) -> str:
    line = (text or "").strip().replace("\r", " ")
    line = line.split("\n")[0].strip()
    return line[:limit] + ("…" if len(line) > limit else "")


def _flatten(text: str, limit: int = 220) -> str:
    """Collapse all whitespace/newlines to single spaces and truncate."""
    t = " ".join((text or "").split())
    return t[:limit] + ("…" if len(t) > limit else "")


_BULLET_RE = re.compile(r"^\s*(?:[-*•●▪‣⁃o]|\d+[.)])\s+")


def _summarize_description(text: str, limit: int = 320,
                           max_words: int | None = None) -> str:
    """Summarise a description. If it's bulleted/multi-line, join every point
    (not just the first) into one '; '-separated line.

    ``max_words`` caps the result to a word count (used for epics).
    """
    raw = (text or "").strip()
    if not raw:
        return ""
    lines = [_BULLET_RE.sub("", ln).strip() for ln in raw.splitlines()]
    points = [ln for ln in lines if ln]
    if not points:
        return ""
    s = points[0] if len(points) == 1 else "; ".join(points)
    if max_words:
        words = s.split()
        if len(words) > max_words:
            return " ".join(words[:max_words]) + "…"
        return s
    return s[:limit] + ("…" if len(s) > limit else "")


def _in_month(value, year: int, month: int) -> bool:
    d = fields_mod.to_date(value)
    return bool(d and d.year == year and d.month == month)


def _build_row(client: JiraClient, local_fields: LocalFieldsStore, s: dict,
               tgt: str, today) -> dict:
    """Build one sub-task row (RAG status + recent comments)."""
    sf = s["fields"]
    status_name = (sf.get("status") or {}).get("name", "")
    cat = ((sf.get("status") or {}).get("statusCategory") or {}).get("key", "")
    end = sf.get(tgt)
    cmts = client.recent_comments(s["key"])
    label, colour = rag_status(status_name, cat, end, cmts["allText"], today)
    updates = [
        {"text": _flatten(c["text"], 180), "author": c["author"], "date": c["date"]}
        for c in cmts["recent"] if c["text"].strip()
    ]
    lf = local_fields.get(s["key"])
    return {
        "key": s["key"],
        "summary": sf.get("summary", ""),
        "start": str(sf.get(START_DATE_FIELD) or "")[:10] or "—",
        "end": str(end or "")[:10] or "—",
        "status": status_name,
        "rag": label,
        "ragColor": "#%02X%02X%02X" % (colour[0], colour[1], colour[2]),
        "_color": colour,
        "who": (sf.get("assignee") or {}).get("displayName", "Unassigned"),
        "reporter": (sf.get("reporter") or {}).get("displayName", "\u2014"),
        "component": (lf.get("component") or {}).get("name", ""),
        "developer": (lf.get("developer") or {}).get("displayName", ""),
        "assignedGroup": (lf.get("assignedGroup") or {}).get("name", ""),
        "updates": updates,
        "commentsFallback": cmts.get("fallback", False),
    }


def gather(client: JiraClient, year: int, month: int,
           local_fields: LocalFieldsStore, assignee: str | None = None) -> list[dict]:
    """Collect epics -> stories -> in-month sub-tasks with RAG status.

    A story (and its epic) is included if ANY of these hold:
      * the epic is assigned to the user, or
      * the story itself is assigned to or reported by the user, or
      * the story holds a sub-task assigned to or reported by the user.
    In every case the story is only shown when it has at least one sub-task
    whose Target Completion Date falls in the selected month; all such
    sub-tasks are then listed regardless of who they are assigned to.
    ``assignee`` is an accountId; defaults to the current user.
    """
    today = date.today()
    tgt = fields_mod.TARGET_COMPLETION_FIELD
    who = "currentUser()" if not assignee else f'"{assignee}"'

    def chunks(seq, n=50):
        for i in range(0, len(seq), n):
            yield seq[i:i + n]

    # 1. Relevant "story" (standard-issue) keys.
    story_keys: set[str] = set()

    #    (a) stories under epics assigned to the user
    epics_mine = client.search(
        f"assignee = {who} AND issuetype = Epic AND statusCategory != Done",
        fields=["summary"])
    for batch in chunks([e["key"] for e in epics_mine]):
        keys = ", ".join(f'"{k}"' for k in batch)
        for ch in client.search(f"parent in ({keys})", fields=["issuetype"]):
            if not (ch["fields"].get("issuetype") or {}).get("subtask"):
                story_keys.add(ch["key"])

    #    (b) stories that hold a sub-task assigned to / reported by the user
    #        (with an in-month target).
    for s in client.search(
            f"(assignee = {who} OR reporter = {who}) AND parent is not EMPTY",
            fields=["parent", "issuetype", tgt]):
        if not (s["fields"].get("issuetype") or {}).get("subtask"):
            continue
        if _in_month(s["fields"].get(tgt), year, month):
            parent = (s["fields"].get("parent") or {}).get("key")
            if parent:
                story_keys.add(parent)

    #    (c) stories assigned to / reported by the user directly (epic/sub-tasks
    #        may not be). Their in-month sub-tasks should still surface.
    for st in client.search(
            f"(assignee = {who} OR reporter = {who}) AND issuetype != Epic",
            fields=["issuetype"]):
        if not (st["fields"].get("issuetype") or {}).get("subtask"):
            story_keys.add(st["key"])

    if not story_keys:
        return []

    # 2. Story details (summary, description, parent epic).
    stories: dict[str, dict] = {}
    for batch in chunks(list(story_keys)):
        keys = ", ".join(f'"{k}"' for k in batch)
        for st in client.search(
                f"key in ({keys})",
                fields=["summary", "description", "parent", "issuetype"]):
            if not (st["fields"].get("issuetype") or {}).get("subtask"):
                stories[st["key"]] = st

    # 3. Epic details for the epics those stories belong to.
    epic_keys = {(st["fields"].get("parent") or {}).get("key")
                 for st in stories.values()}
    epic_keys.discard(None)
    epics: dict[str, dict] = {}
    for batch in chunks(list(epic_keys)):
        keys = ", ".join(f'"{k}"' for k in batch)
        for e in client.search(f"key in ({keys})",
                               fields=["summary", "description"]):
            epics[e["key"]] = e

    # 4. Build epic -> story -> in-month sub-task rows.
    grouped: dict[str, dict] = {}
    for skey, st in stories.items():
        sf = st["fields"]
        subs = client.search(
            f'parent = "{skey}" ORDER BY summary',
            fields=["summary", "status", "assignee", "reporter",
                    START_DATE_FIELD, tgt])
        month_subs = [s for s in subs if _in_month(s["fields"].get(tgt), year, month)]
        if not month_subs:
            continue
        total = len(subs)
        done = sum(1 for s in subs
                   if ((s["fields"].get("status") or {}).get("statusCategory") or {})
                   .get("key") == "done")
        story_out = {
            "key": skey,
            "summary": sf.get("summary", ""),
            "description": _summarize_description(adf_to_text(sf.get("description"))),
            "progress": f"{done}/{total} sub-tasks done",
            "subs": [_build_row(client, local_fields, s, tgt, today) for s in month_subs],
        }

        ep_key = (sf.get("parent") or {}).get("key")
        if ep_key and ep_key in epics:
            ef = epics[ep_key]["fields"]
            grp = grouped.setdefault(ep_key, {
                "key": ep_key,
                "summary": ef.get("summary", ""),
                "description": _summarize_description(
                    adf_to_text(ef.get("description")), max_words=40),
                "stories": [],
            })
        else:
            grp = grouped.setdefault("", {
                "key": "", "summary": "(No epic)", "description": "",
                "stories": [],
            })
        grp["stories"].append(story_out)

    result = list(grouped.values())
    for ep in result:
        ep["stories"].sort(key=lambda s: s["key"])
    result.sort(key=lambda e: e["key"])
    return result


def month_label(year: int, month: int) -> str:
    return f"{calendar.month_name[month]} {year}"


# --- PowerPoint rendering --------------------------------------------------

COLS = ["Story / Sub-task", "Component", "Developer", "Assigned Group",
        "Start", "Target end", "Status", "Progress", "Responsible",
        "Reported by", "Recent comments (4 wks)"]
COL_WIDTHS = [2.6, 0.9, 1.0, 1.0, 0.7, 0.7, 0.8, 0.9, 1.0, 1.0, 1.7]  # ~12.3 on 13.33 slide


def _set_cell(cell, runs, *, bold=False, size=9, color=DARK, fill=None,
              align=PP_ALIGN.LEFT):
    """runs: list of (text, bold) lines, or a plain string."""
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Inches(0.05)
    cell.margin_right = Inches(0.05)
    cell.margin_top = Inches(0.02)
    cell.margin_bottom = Inches(0.02)
    if fill is not None:
        cell.fill.solid()
        cell.fill.fore_color.rgb = fill
    tf = cell.text_frame
    tf.word_wrap = True
    lines = runs if isinstance(runs, list) else [(runs, bold)]
    for i, (text, b) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.size = Pt(size)
        r.font.bold = b
        r.font.color.rgb = color


def _add_title(slide, text: str) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.7))
    tp = title_box.text_frame.paragraphs[0]
    tp.word_wrap = True
    run = tp.add_run()
    run.text = text
    run.font.size = Pt(24)
    run.font.bold = True
    run.font.color.rgb = HEADER_BG


def _render_epic_slide(prs, ep: dict, year: int, month: int,
                       title_suffix: str) -> None:
    """Render a single epic onto its own slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    etitle = f'{ep["key"]}: {ep["summary"]}' if ep["key"] else ep["summary"]
    _add_title(slide, f"{month_label(year, month)} – {etitle}{title_suffix}")

    if ep["description"]:
        desc_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.9), Inches(12.5), Inches(0.5))
        dtf = desc_box.text_frame
        dtf.word_wrap = True
        dr = dtf.paragraphs[0].add_run()
        dr.text = ep["description"]
        dr.font.size = Pt(11)
        dr.font.color.rgb = DARK

    table_top = 1.45 if ep["description"] else 1.05

    # Flatten this epic's stories -> sub-tasks into table rows.
    flat = []
    for st in ep["stories"]:
        s_first = True
        for sub in st["subs"]:
            flat.append({"story": st, "sub": sub, "s_first": s_first})
            s_first = False

    if not flat:
        return

    n_rows = len(flat) + 1  # + header
    rows_h = min(0.55, max(0.28, 5.8 / len(flat)))
    table_h = Inches(0.4 + rows_h * len(flat))
    gf = slide.shapes.add_table(n_rows, len(COLS), Inches(0.4), Inches(table_top),
                                Inches(sum(COL_WIDTHS)), table_h)
    table = gf.table
    table.first_row = True
    for i, w in enumerate(COL_WIDTHS):
        table.columns[i].width = Inches(w)

    # Header
    for c, name in enumerate(COLS):
        _set_cell(table.cell(0, c), name, bold=True, size=10, color=WHITE,
                  fill=HEADER_BG, align=PP_ALIGN.CENTER)

    # Body
    for idx, row in enumerate(flat):
        r = idx + 1
        st, sub = row["story"], row["sub"]
        # Story / sub-task cell: story header (first row of story) + sub-task line
        story_lines = []
        if row["s_first"]:
            story_lines.append((f'{st["key"]}: {st["summary"]}  ({st["progress"]})', True))
            if st.get("description"):
                story_lines.append((st["description"], False))
        story_lines.append((f'↳ {sub["key"]}: {sub["summary"]}', False))
        _set_cell(table.cell(r, 0), story_lines, size=9)
        _set_cell(table.cell(r, 1), sub.get("component", ""), size=8)
        _set_cell(table.cell(r, 2), sub.get("developer", ""), size=8)
        _set_cell(table.cell(r, 3), sub.get("assignedGroup", ""), size=8)
        _set_cell(table.cell(r, 4), sub["start"], size=8, align=PP_ALIGN.CENTER)
        _set_cell(table.cell(r, 5), sub["end"], size=8, align=PP_ALIGN.CENTER)
        _set_cell(table.cell(r, 6), [(sub["rag"], True)], size=8, color=WHITE,
                  fill=sub["_color"], align=PP_ALIGN.CENTER)
        _set_cell(table.cell(r, 7), sub.get("status", ""), size=8,
                  align=PP_ALIGN.CENTER)
        _set_cell(table.cell(r, 8), sub["who"], size=8)
        _set_cell(table.cell(r, 9), sub.get("reporter", "\u2014"), size=8)
        updates = sub.get("updates") or []
        if updates:
            fb = " (latest, none in 4 wks)" if sub.get("commentsFallback") else ""
            lines = []
            for u in updates:
                lines.append((u["text"], False))
                lines.append(("- {0}, {1}{2}".format(u["author"], u["date"], fb), False))
            _set_cell(table.cell(r, 10), lines, size=7, color=DARK)
        else:
            _set_cell(table.cell(r, 10), [("No comments", False)], size=7, color=GREY)


def build_pptx(epics: list[dict], year: int, month: int,
               owner: str | None = None) -> bytes:
    title_suffix = f"  —  {owner}" if owner else ""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    rendered = [ep for ep in epics if any(st["subs"] for st in ep["stories"])]

    if not rendered:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        _add_title(slide, f"Monthly Report – {month_label(year, month)}{title_suffix}")
        box = slide.shapes.add_textbox(Inches(0.4), Inches(1.2), Inches(12), Inches(1))
        box.text_frame.paragraphs[0].add_run().text = (
            f"No sub-tasks with a Target Completion Date in {month_label(year, month)}."
        )
        out = io.BytesIO(); prs.save(out); return out.getvalue()

    # One epic per slide.
    for ep in rendered:
        _render_epic_slide(prs, ep, year, month, title_suffix)

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
